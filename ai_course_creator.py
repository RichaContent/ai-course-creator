import streamlit as st
import os
from openai import OpenAI
from pptx import Presentation
from pptx.util import Pt
from io import BytesIO
from docx import Document
import tempfile
import PyPDF2
import docx2txt
import zipfile

# Initialize OpenAI
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

st.set_page_config(page_title="AI Course Creator", layout="centered")
st.title("📚 AI Course Creator for Trainers")

# Built-in prompt templates
templates = {
    "Leadership": """You are an expert instructional designer.

Create a delivery-ready, 90-minute training course on "Effective Leadership for First-Time Managers" for an audience of new managers.

Requirements:
✅ Tabular Course Outline (Time | Activity | Description).
✅ Detailed Facilitator Guide with practical instructions, public domain definitions, examples.
✅ Participant Workbook with reflective activities and exercises.
✅ Quiz (5-8 questions with answers).
✅ Slide Deck with clear titles and bullet points.
✅ Public domain references only, no placeholders.
✅ Practical examples, real scenarios.
✅ Tone: Professional yet engaging.

If notes or feedback are provided, integrate them while maintaining consistency.

Return sections:
## COURSE_OUTLINE
## FACILITATOR_GUIDE
## PARTICIPANT_WORKBOOK
## QUIZ
## SLIDE_DECK
""",
    "Time Management": """You are an expert instructional designer.

Create a delivery-ready, 60-minute training course on "Time Management for Busy Professionals."

Requirements:
✅ Tabular Course Outline.
✅ Facilitator Guide with detailed talking points, examples.
✅ Participant Workbook with practical activities (prioritization, Eisenhower Matrix).
✅ Quiz with 5-8 aligned questions and answers.
✅ Slide Deck with actionable tips, examples, and reflection prompts.
✅ Clear, jargon-free, public domain references only.
✅ Practical case studies included.

If notes or feedback are provided, integrate them seamlessly.

Return sections:
## COURSE_OUTLINE
## FACILITATOR_GUIDE
## PARTICIPANT_WORKBOOK
## QUIZ
## SLIDE_DECK
""",
    "Feedback & Performance Management": """You are an expert instructional designer.

Create a delivery-ready, 90-minute training course on "Feedback and Performance Management for Middle Managers."

Requirements:
✅ Tabular Course Outline.
✅ Detailed Facilitator Guide with sample dialogues, scenarios, SBIS model explanation.
✅ Participant Workbook with exercises on preparing feedback, role-plays.
✅ Quiz with practical scenario-based questions and answers.
✅ Slide Deck aligned to the course flow.
✅ All examples and quotes should be from public domain or generalized.

If notes or uploaded files are provided, incorporate them while maintaining previous structure.

Return sections:
## COURSE_OUTLINE
## FACILITATOR_GUIDE
## PARTICIPANT_WORKBOOK
## QUIZ
## SLIDE_DECK
""",
    "Emotional Intelligence": """You are an expert instructional designer.

Create a delivery-ready, 90-minute training course on "Emotional Intelligence in the Workplace."

Requirements:
✅ Tabular Course Outline.
✅ Facilitator Guide with explanations of self-awareness, self-regulation, empathy, etc.
✅ Participant Workbook with reflection logs, activities.
✅ Quiz aligned with objectives, with answers.
✅ Slide Deck with clear visuals, definitions, and practical applications.
✅ Public domain references, no placeholders.

If notes or feedback are provided, incorporate while retaining structure.

Return sections:
## COURSE_OUTLINE
## FACILITATOR_GUIDE
## PARTICIPANT_WORKBOOK
## QUIZ
## SLIDE_DECK
"""
}

# User selects template
st.header("Step 1: Select Training Topic")
selected_template = st.selectbox("Choose a Course Template", list(templates.keys()))

# Optionally edit topic/audience/duration
topic = st.text_input("Course Topic (Edit if needed)", selected_template)
audience = st.text_input("Target Audience", "Middle Management")
duration = st.slider("Duration (minutes)", 30, 240, 90, step=15)
tone = st.selectbox("Tone of Voice", ["Professional", "Conversational", "Inspirational"])

# Step 2: Upload files and notes
st.header("Step 2: Upload Reference Files & Notes (Optional)")
uploaded_files = st.file_uploader("Upload PDF, DOCX, PPTX files", type=["pdf", "docx", "pptx"], accept_multiple_files=True)
notes = st.text_area("Notes or Specific Instructions (Optional)")

# Step 3: Feedback for revision
st.header("Step 3: Feedback for Revision (Optional)")
feedback = st.text_area("Feedback for refinement (Optional)")

# File extraction
def extract_uploaded_text(files):
    text = ""
    for file in files:
        if file.name.endswith(".pdf"):
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        elif file.name.endswith(".docx"):
            text += docx2txt.process(file) + "\n"
        elif file.name.endswith(".pptx"):
            ppt = Presentation(file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    return text[:8000]  # limit for token safety

# Slide deck generation
def generate_slide_deck(slide_content):
    prs = Presentation()
    for block in slide_content.strip().split("\n\n"):
        lines = block.strip().split("\n")
        if not lines:
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = lines[0]
        textbox = slide.placeholders[1]
        for bullet in lines[1:]:
            p = textbox.text_frame.add_paragraph()
            p.text = bullet
            p.font.size = Pt(20)
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

# DOCX generation
def save_docx(content, filename):
    doc = Document()
    for line in content.strip().split("\n"):
        doc.add_paragraph(line)
    temp_path = os.path.join(tempfile.gettempdir(), filename)
    doc.save(temp_path)
    return temp_path

# Generate button
if st.button("🚀 Generate Course Materials"):
    with st.spinner("Generating your comprehensive course materials..."):

        extracted_text = extract_uploaded_text(uploaded_files) if uploaded_files else ""
        ref_part = f"\nReference Content:\n{extracted_text}\n" if extracted_text else ""
        notes_part = f"\nNotes:\n{notes}\n" if notes else ""
        feedback_part = f"\nFeedback for refinement:\n{feedback}\n" if feedback else ""

        prompt = templates[selected_template] + ref_part + notes_part + feedback_part

        try:
            completion = client.chat.completions.create(
                model="gpt-4o-preview",
                messages=[
                    {"role": "system", "content": "You are a precise instructional designer generating complete, ready-to-use materials for corporate trainers."},
                    {"role": "user", "content": prompt}
                ]
            )
            response_content = completion.choices[0].message.content

            sections = {key: "" for key in ["COURSE_OUTLINE", "FACILITATOR_GUIDE", "PARTICIPANT_WORKBOOK", "QUIZ", "SLIDE_DECK"]}
            current_section = None
            for line in response_content.splitlines():
                line = line.strip()
                if line.startswith("## ") and line[3:] in sections:
                    current_section = line[3:]
                elif current_section:
                    sections[current_section] += line + "\n"

            # Save DOCX
            outline_file = save_docx(sections["COURSE_OUTLINE"], "Course_Outline.docx")
            guide_file = save_docx(sections["FACILITATOR_GUIDE"], "Facilitator_Guide.docx")
            workbook_file = save_docx(sections["PARTICIPANT_WORKBOOK"], "Participant_Workbook.docx")
            quiz_file = save_docx(sections["QUIZ"], "Quiz.docx")
            slide_deck_file = generate_slide_deck(sections["SLIDE_DECK"])

            # Create ZIP
            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, "w") as zf:
                zf.write(outline_file, "Course_Outline.docx")
                zf.write(guide_file, "Facilitator_Guide.docx")
                zf.write(workbook_file, "Participant_Workbook.docx")
                zf.write(quiz_file, "Quiz.docx")
                zf.writestr("Slide_Deck.pptx", slide_deck_file.getvalue())
            zip_buffer.seek(0)

            # Downloads
            st.success("✅ All materials generated successfully!")
            st.download_button("📥 Download All as ZIP", data=zip_buffer, file_name="Course_Materials.zip")
            st.download_button("Download Course Outline", open(outline_file, "rb"), file_name="Course_Outline.docx")
            st.download_button("Download Facilitator Guide", open(guide_file, "rb"), file_name="Facilitator_Guide.docx")
            st.download_button("Download Workbook", open(workbook_file, "rb"), file_name="Participant_Workbook.docx")
            st.download_button("Download Quiz", open(quiz_file, "rb"), file_name="Quiz.docx")
            st.download_button("Download Slide Deck", slide_deck_file, file_name="Slide_Deck.pptx")

            # Token tracking
            try:
                tokens_used = completion.usage.total_tokens
                cost_estimate = round(tokens_used / 1000 * 0.03, 4)
                st.caption(f"Used {tokens_used} tokens · Estimated cost: ${cost_estimate:.4f}")
            except:
                pass

        except Exception as e:
            st.error(f"❌ Error generating content: {e}")
